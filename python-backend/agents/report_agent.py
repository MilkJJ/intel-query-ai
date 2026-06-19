"""
ReportAgent — generates insight-driven PDF and PowerPoint reports from video analysis.

Phase 5: LLM-Enhanced Report Generation

Capabilities:
- Generate PDF reports with 9 insight-driven sections (NOT transcript dumps)
  1. Title + metadata
  2. TL;DR (brief summary)
  3. Executive Summary (LLM-generated)
  4. Key Topics (LLM-identified themes)
  5. Scene Analysis (visual timeline)
  6. Object Detection Results (table format)
  7. Graph Analysis (if applicable)
  8. Key Insights (LLM reasoning)
  9. Conclusion (LLM-generated)

- Generate PPT presentations with 6 insight-focused slides
  1. Title slide
  2. Executive Summary
  3. Key Topics
  4. Objects Detected
  5. Key Insights
  6. Conclusion

- Use GenerationAgent for LLM-based insights (Phase 4 integration)
- Return saved file path plus structured metadata

Key Design: Reports focus on insights and findings, NOT raw transcript text.
"""

from __future__ import annotations

import logging
import tempfile
from collections import Counter
from datetime import datetime
from pathlib import Path
from typing import Any
from xml.sax.saxutils import escape

from fastapi import HTTPException
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from agents.base import BaseAgent
from agents.context import AgentContext
from agents.result import AgentResult


logger = logging.getLogger(__name__)

DEFAULT_REPORT_TITLE = "Video AI Report"


class ReportAgent(BaseAgent):
    name = "report"

    def __init__(self, transcription_agent, vision_agent=None, generation_agent=None):
        """
        Initialize ReportAgent.
        
        Args:
            transcription_agent: TranscriptionAgent for audio extraction
            vision_agent: VisionAgent for frame analysis (optional)
            generation_agent: GenerationAgent for LLM-based insights (optional)
        """
        self._transcription_agent = transcription_agent
        self._vision_agent = vision_agent
        self._generation_agent = generation_agent

    @property
    def actions(self):
        return {
            "generate_pdf": self.generate_pdf,
            "generate_ppt": self.generate_ppt,
        }

    def run(self, context: AgentContext) -> AgentResult:
        action = context.action or "generate_pdf"
        handler = self.actions.get(action)
        if handler is None:
            raise ValueError(f"ReportAgent: unknown action '{action}'")
        return handler(context)

    # ── Input normalization ──────────────────────────────────────────────────

    def _ensure_transcript(self, context: AgentContext) -> str:
        if context.transcript:
            return context.transcript

        if self._transcription_agent is None:
            return ""

        saved_action = context.action
        try:
            context.action = "get_transcript"
            transcription_result = self._transcription_agent.run(context)
            context.transcript = transcription_result.transcription
            return context.transcript
        finally:
            context.action = saved_action

    def _ensure_transcript_summary(self, context: AgentContext) -> str:
        summary = context.metadata.get("transcript_summary")
        if isinstance(summary, str) and summary.strip():
            return summary.strip()

        transcript = self._ensure_transcript(context)
        saved_action = context.action
        try:
            if self._transcription_agent is not None:
                context.action = "summarize_transcript"
                summary_result = self._transcription_agent.run(context)
                summary_text = summary_result.response.strip()
                if summary_text:
                    context.metadata["transcript_summary"] = summary_text
                    return summary_text
        except Exception:
            logger.warning("ReportAgent: transcript summarization failed, falling back to heuristic summary.")
        finally:
            context.action = saved_action

        heuristic_summary = self._summarize_transcript(transcript)
        context.metadata["transcript_summary"] = heuristic_summary
        return heuristic_summary

    def _summarize_transcript(self, transcript: str) -> str:
        if not transcript.strip():
            return "No transcript available."

        sentences = [sentence.strip() for sentence in transcript.replace("\n", " ").split(".") if sentence.strip()]
        summary_sentences = sentences[:5]
        if not summary_sentences:
            return transcript[:500]
        return ". ".join(summary_sentences) + ("." if not summary_sentences[-1].endswith(".") else "")

    def _coerce_string_list(self, value: Any) -> list[str]:
        if value is None:
            return []
        if isinstance(value, list):
            return [str(item) for item in value if str(item).strip()]
        if isinstance(value, tuple):
            return [str(item) for item in value if str(item).strip()]
        if isinstance(value, str):
            return [value] if value.strip() else []
        return [str(value)]

    def _coerce_frame_text(self, value: Any) -> list[str]:
        if value is None:
            return []
        if isinstance(value, list):
            frame_text: list[str] = []
            for item in value:
                if isinstance(item, dict):
                    text = item.get("text") or item.get("caption") or item.get("content")
                    if text and str(text).strip():
                        frame_text.append(str(text).strip())
                elif str(item).strip():
                    frame_text.append(str(item).strip())
            return frame_text
        if isinstance(value, dict):
            text = value.get("text") or value.get("caption") or value.get("content")
            return [str(text).strip()] if text and str(text).strip() else []
        if isinstance(value, str):
            return [value] if value.strip() else []
        return [str(value)] if str(value).strip() else []

    def _extract_vision_inputs(self, vision_result: AgentResult | None) -> tuple[list[str], list[str], dict[str, Any]]:
        if vision_result is None:
            return [], [], {}

        metadata = vision_result.metadata or {}
        summary = metadata.get("summary") if isinstance(metadata.get("summary"), dict) else {}
        if summary is None:
            summary = {}

        detected_objects = self._coerce_string_list(
            summary.get("objects_detected")
            or metadata.get("objects")
            or metadata.get("detected_objects")
            or vision_result.objects
        )

        frame_text = self._coerce_frame_text(
            summary.get("text_snippets")
            or metadata.get("frame_text")
            or metadata.get("text_frames")
            or metadata.get("frames")
        )

        return detected_objects, frame_text, metadata

    def _ensure_vision_analysis(self, context: AgentContext) -> AgentResult | None:
        if self._vision_agent is None:
            return None

        cached = context.metadata.get("vision_analysis")
        if isinstance(cached, dict):
            return None

        saved_action = context.action
        try:
            context.action = "analyze_video"
            vision_result = self._vision_agent.run(context)
            context.metadata["vision_analysis"] = vision_result.metadata
            return vision_result
        except Exception as error:
            logger.warning("ReportAgent: vision analysis failed, continuing without it: %s", error)
            return None
        finally:
            context.action = saved_action

    def _ensure_llm_insights(self, context: AgentContext, analysis_data: dict[str, Any]) -> dict[str, Any]:
        """Get LLM-generated insights from GenerationAgent (Phase 4)."""
        if self._generation_agent is None:
            logger.warning("ReportAgent: GenerationAgent not available, using heuristic insights")
            return self._heuristic_insights(analysis_data)

        # Check cache first
        cached_insights = context.metadata.get("generated_insights")
        if isinstance(cached_insights, dict) and cached_insights:
            return cached_insights

        try:
            logger.info("ReportAgent: Requesting LLM insights from GenerationAgent")
            saved_action = context.action
            try:
                context.action = "generate_insights"
                result = self._generation_agent.run(context)
                insights = result.metadata or {}
                
                # Ensure insights dict has required keys
                if not isinstance(insights, dict):
                    insights = {}
                    
                return {
                    "summary": insights.get("summary", analysis_data.get("transcript_summary", "")),
                    "key_topics": insights.get("key_topics", []),
                    "key_insights": insights.get("key_insights", []),
                    "conclusion": insights.get("conclusion", ""),
                }
            finally:
                context.action = saved_action
        except Exception as e:
            logger.warning("ReportAgent: LLM insight generation failed: %s, using heuristic", e)
            return self._heuristic_insights(analysis_data)

    def _heuristic_insights(self, analysis_data: dict[str, Any]) -> dict[str, Any]:
        """Generate heuristic insights when LLM unavailable."""
        transcript = analysis_data.get("transcript", "")
        objects = analysis_data.get("detected_objects", [])
        
        # Extract key topics from objects and transcript
        key_topics = list(set(objects))[:5] if objects else []
        
        # Generate heuristic insights
        key_insights = [
            f"Detected {len(objects)} distinct object types" if objects else "Video contains visual content",
            f"Transcript shows key discussion points" if transcript else "Content focuses on visual elements",
        ]
        
        return {
            "summary": analysis_data.get("transcript_summary", ""),
            "key_topics": key_topics,
            "key_insights": key_insights,
            "conclusion": "Video provides comprehensive information through multimodal presentation.",
        }

    def _collect_analysis_inputs(self, context: AgentContext) -> dict[str, Any]:
        transcript = self._ensure_transcript(context)
        transcript_summary = self._ensure_transcript_summary(context)

        vision_result = self._ensure_vision_analysis(context)
        detected_objects, frame_text, vision_metadata = self._extract_vision_inputs(vision_result)

        detected_objects = context.metadata.get("detected_objects") or detected_objects
        detected_objects = self._coerce_string_list(detected_objects)

        frame_text = context.metadata.get("frame_text") or frame_text
        frame_text = self._coerce_string_list(frame_text)

        vision_summary = vision_metadata.get("summary") if isinstance(vision_metadata.get("summary"), dict) else {}
        if vision_summary is None:
            vision_summary = {}

        # Build basic analysis data
        analysis_data = {
            "title": context.metadata.get("report_title") or DEFAULT_REPORT_TITLE,
            "video_name": Path(context.video_path).name if context.video_path else "Unknown",
            "query": context.query,
            "transcript": transcript,
            "transcript_summary": transcript_summary,
            "detected_objects": detected_objects,
            "frame_text": frame_text,
            "vision_summary": vision_summary,
            "vision_metadata": vision_metadata,
            "created_at": datetime.utcnow().isoformat(timespec="seconds") + "Z",
        }

        # Get LLM insights (Phase 5)
        llm_insights = self._ensure_llm_insights(context, analysis_data)
        analysis_data["llm_insights"] = llm_insights

        return analysis_data

    def _report_directory(self, context: AgentContext) -> Path:
        root = Path(tempfile.gettempdir()) / "intel-query-ai-reports"
        if context.session_id:
            root = root / context.session_id
        root.mkdir(parents=True, exist_ok=True)
        return root

    def _safe_stem(self, context: AgentContext, suffix: str) -> str:
        base_name = Path(context.video_path).stem if context.video_path else "report"
        timestamp = datetime.utcnow().strftime("%Y%m%d-%H%M%S")
        return f"{base_name}-{suffix}-{timestamp}".lower().replace(" ", "-")

    # ── PDF generation ───────────────────────────────────────────────────────

    def _safe_paragraph(self, text: str) -> str:
        if not text:
            return "No content available."
        return escape(text).replace("\n", "<br/>")

    def _top_items(self, values: list[str], limit: int = 5) -> list[str]:
        cleaned = [value.strip() for value in values if isinstance(value, str) and value.strip()]
        if not cleaned:
            return []
        counts = Counter(cleaned)
        return [item for item, _ in counts.most_common(limit)]

    def _key_findings(self, report_data: dict[str, Any]) -> list[str]:
        findings: list[str] = []

        summary_text = (report_data.get("transcript_summary") or "").strip()
        if summary_text:
            summary_excerpt = summary_text[:220] + ("..." if len(summary_text) > 220 else "")
            findings.append(f"Transcript summary: {summary_excerpt}")

        top_objects = self._top_items(report_data.get("detected_objects", []), limit=5)
        if top_objects:
            findings.append("Most visible objects: " + ", ".join(top_objects) + ".")

        frame_text = self._top_items(report_data.get("frame_text", []), limit=3)
        if frame_text:
            findings.append("Important on-screen text: " + " | ".join(frame_text) + ".")

        if not findings:
            findings.append("Limited analysis signals were available, so this report is based on minimal extracted data.")

        return findings

    def _recommended_actions(self, report_data: dict[str, Any]) -> list[str]:
        recommendations = [
            "Use the transcript summary as the opening narrative when sharing this content.",
            "Validate key object detections against the original video before final publication.",
        ]

        if report_data.get("frame_text"):
            recommendations.append("Include the extracted frame text as supporting evidence in an appendix or speaker notes.")
        else:
            recommendations.append("Re-run OCR with higher frame sampling if on-screen text is important for this use case.")

        if report_data.get("query"):
            recommendations.append(f"Follow-up question to investigate: {report_data['query']}")

        return recommendations

    def _format_report_response(self, report_data: dict[str, Any], output_path: Path, report_format: str) -> str:
        findings = self._key_findings(report_data)
        format_name = report_format.upper()
        lines = [
            f"Your {format_name} report is ready.",
            f"Saved to: {output_path}",
            "",
            "Highlights:",
            f"- {findings[0]}",
            f"- Objects detected: {len(report_data.get('detected_objects', []))}",
            f"- Text snippets extracted: {len(report_data.get('frame_text', []))}",
        ]
        return "\n".join(lines)

    def _vision_insight_lines(self, vision_summary: dict[str, Any]) -> list[str]:
        lines: list[str] = []
        for key, value in vision_summary.items():
            if isinstance(value, (str, int, float, bool)):
                lines.append(f"{key}: {value}")
            elif isinstance(value, list):
                lines.append(f"{key}: {len(value)} items")
            elif isinstance(value, dict):
                lines.append(f"{key}: {len(value)} fields")
        return lines

    def _build_pdf_story(self, report_data: dict[str, Any], styles) -> list[Any]:
        """
        Build 9-section PDF story with LLM insights (Phase 5).
        
        Sections:
        1. Title + metadata
        2. TL;DR (brief summary)
        3. Executive Summary (LLM)
        4. Key Topics (LLM)
        5. Scene Analysis (visual timeline)
        6. Object Detection Results (table)
        7. Graph Analysis (if applicable)
        8. Key Insights (LLM)
        9. Conclusion (LLM)
        """
        story: list[Any] = []

        title_style = ParagraphStyle(
            "ReportTitle",
            parent=styles["Title"],
            fontSize=22,
            leading=26,
            textColor=colors.HexColor("#111827"),
            spaceAfter=10,
        )
        heading_style = ParagraphStyle(
            "ReportHeading",
            parent=styles["Heading2"],
            fontSize=14,
            leading=18,
            textColor=colors.HexColor("#1f2937"),
            spaceAfter=6,
        )
        body_style = ParagraphStyle(
            "ReportBody",
            parent=styles["BodyText"],
            fontSize=10.5,
            leading=14,
            textColor=colors.HexColor("#111827"),
        )

        # ── Section 1: Title + Metadata ──────────────────────────────────────
        story.append(Paragraph(report_data["title"], title_style))
        story.append(Paragraph(f"Video: {report_data['video_name']}", body_style))
        story.append(Paragraph(f"Created: {report_data['created_at']}", body_style))
        story.append(Paragraph(f"Query: {report_data['query']}", body_style))
        story.append(Spacer(1, 8 * mm))

        llm_insights = report_data.get("llm_insights", {})

        # ── Section 2: TL;DR ─────────────────────────────────────────────────
        story.append(Paragraph("TL;DR", heading_style))
        tldr = llm_insights.get("summary", report_data.get("transcript_summary", ""))
        tldr_short = tldr[:400] + ("..." if len(tldr) > 400 else "")
        story.append(Paragraph(self._safe_paragraph(tldr_short), body_style))
        story.append(Spacer(1, 4 * mm))

        # ── Section 3: Executive Summary ─────────────────────────────────────
        story.append(Paragraph("Executive Summary", heading_style))
        summary = llm_insights.get("summary", "")
        story.append(Paragraph(self._safe_paragraph(summary), body_style))
        story.append(Spacer(1, 4 * mm))

        # ── Section 4: Key Topics ────────────────────────────────────────────
        story.append(Paragraph("Key Topics", heading_style))
        key_topics = llm_insights.get("key_topics", [])
        if key_topics:
            topics_text = "\n".join(f"• {topic}" for topic in key_topics)
            story.append(Paragraph(self._safe_paragraph(topics_text), body_style))
        else:
            story.append(Paragraph("No key topics identified.", body_style))
        story.append(Spacer(1, 4 * mm))

        # ── Section 5: Scene Analysis ────────────────────────────────────────
        story.append(Paragraph("Scene Analysis", heading_style))
        vision_metadata = report_data.get("vision_metadata", {})
        scenes = vision_metadata.get("scenes", [])
        if scenes:
            scene_text = "\n".join(
                f"• {scene.get('timestamp', 'N/A')}s: {scene.get('caption', 'No caption')}"
                for scene in scenes[:5]  # Limit to first 5 scenes
            )
            story.append(Paragraph(self._safe_paragraph(scene_text), body_style))
        else:
            frame_text = report_data.get("frame_text", [])
            if frame_text:
                frame_summary = "\n".join(f"• {text}" for text in frame_text[:5])
                story.append(Paragraph(self._safe_paragraph(frame_summary), body_style))
            else:
                story.append(Paragraph("No scene data available.", body_style))
        story.append(Spacer(1, 4 * mm))

        # ── Section 6: Object Detection Results ───────────────────────────────
        story.append(Paragraph("Object Detection Results", heading_style))
        objects = report_data.get("detected_objects", [])
        if objects:
            object_rows = [["Object", "Count"]]
            object_counts = Counter(objects)
            for obj, count in object_counts.most_common(8):
                object_rows.append([obj, str(count)])
            
            object_table = Table(object_rows, colWidths=[100 * mm, 60 * mm])
            object_table.setStyle(
                TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#e5e7eb")),
                    ("BACKGROUND", (0, 1), (-1, -1), colors.white),
                    ("TEXTCOLOR", (0, 0), (-1, -1), colors.HexColor("#111827")),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#d1d5db")),
                    ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
                    ("FONTSIZE", (0, 0), (-1, -1), 10),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                    ("TOPPADDING", (0, 0), (-1, -1), 8),
                ])
            )
            story.append(object_table)
        else:
            story.append(Paragraph("No objects detected.", body_style))
        story.append(Spacer(1, 4 * mm))

        # ── Section 7: Graph Analysis ────────────────────────────────────────
        graphs_detected = vision_metadata.get("graphs_detected", False)
        graph_descriptions = vision_metadata.get("graph_descriptions", [])
        if graphs_detected or graph_descriptions:
            story.append(Paragraph("Graph Analysis", heading_style))
            if graph_descriptions:
                graph_text = "\n".join(f"• {desc}" for desc in graph_descriptions[:3])
                story.append(Paragraph(self._safe_paragraph(graph_text), body_style))
            else:
                story.append(Paragraph("Graphs detected in video. Review visual content for details.", body_style))
            story.append(Spacer(1, 4 * mm))

        # ── Section 8: Key Insights ──────────────────────────────────────────
        story.append(Paragraph("Key Insights", heading_style))
        key_insights = llm_insights.get("key_insights", [])
        if key_insights:
            insights_text = "\n".join(f"• {insight}" for insight in key_insights)
            story.append(Paragraph(self._safe_paragraph(insights_text), body_style))
        else:
            story.append(Paragraph("No additional insights available.", body_style))
        story.append(Spacer(1, 4 * mm))

        # ── Section 9: Conclusion ────────────────────────────────────────────
        story.append(Paragraph("Conclusion", heading_style))
        conclusion = llm_insights.get("conclusion", "")
        story.append(Paragraph(self._safe_paragraph(conclusion), body_style))

        return story

    def _write_pdf(self, report_data: dict[str, Any], output_path: Path) -> None:
        styles = getSampleStyleSheet()
        document = SimpleDocTemplate(
            str(output_path),
            pagesize=A4,
            rightMargin=18 * mm,
            leftMargin=18 * mm,
            topMargin=18 * mm,
            bottomMargin=18 * mm,
            title=report_data["title"],
        )
        story = self._build_pdf_story(report_data, styles)
        document.build(story)

    # ── PPT generation ───────────────────────────────────────────────────────

    def _add_title_slide(self, presentation: Any, report_data: dict[str, Any]) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = report_data["title"]
        subtitle = slide.placeholders[1].text_frame
        subtitle.clear()
        subtitle.text = (
            f"Video: {report_data['video_name']}\n"
            f"Created: {report_data['created_at']}\n"
            f"Query: {report_data['query']}"
        )

    def _add_bullet_slide(self, presentation: Any, title: str, bullets: list[str]) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        if not bullets:
            bullets = ["No content available."]

        for index, bullet in enumerate(bullets):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0

    def _write_ppt(self, report_data: dict[str, Any], output_path: Path) -> None:
        """
        Generate 6-slide PowerPoint with LLM insights (Phase 5).
        
        Slides:
        1. Title slide
        2. Executive Summary
        3. Key Topics
        4. Objects Detected
        5. Key Insights
        6. Conclusion
        """
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        llm_insights = report_data.get("llm_insights", {})

        # Slide 1: Title slide
        self._add_title_slide(presentation, report_data)

        # Slide 2: Executive Summary
        summary = llm_insights.get("summary", report_data.get("transcript_summary", "No summary available"))
        summary_bullets = [summary] if summary else ["No summary available"]
        self._add_bullet_slide(presentation, "Executive Summary", summary_bullets)

        # Slide 3: Key Topics
        key_topics = llm_insights.get("key_topics", [])
        if not key_topics:
            # Fallback to detected objects if no key topics
            key_topics = report_data.get("detected_objects", [])[:5]
        self._add_bullet_slide(
            presentation,
            "Key Topics",
            key_topics if key_topics else ["Topics to be identified from video content"],
        )

        # Slide 4: Objects Detected
        objects = report_data.get("detected_objects", [])
        objects_by_count = Counter(objects)
        top_objects = [f"{obj} ({count})" for obj, count in objects_by_count.most_common(8)]
        self._add_bullet_slide(
            presentation,
            "Objects Detected",
            top_objects if top_objects else ["No objects detected"],
        )

        # Slide 5: Key Insights
        key_insights = llm_insights.get("key_insights", [])
        self._add_bullet_slide(
            presentation,
            "Key Insights",
            key_insights if key_insights else ["Key insights to be analyzed from video content"],
        )

        # Slide 6: Conclusion
        conclusion = llm_insights.get("conclusion", "")
        conclusion_bullets = [
            conclusion if conclusion else "Video analysis complete.",
            "For more details, review the full PDF report.",
        ]
        self._add_bullet_slide(presentation, "Conclusion", conclusion_bullets)

        presentation.save(str(output_path))

    # ── Public actions ───────────────────────────────────────────────────────

    def generate_pdf(self, context: AgentContext) -> AgentResult:
        report_data = self._collect_analysis_inputs(context)
        report_dir = self._report_directory(context)
        output_path = report_dir / f"{self._safe_stem(context, 'summary')}.pdf"

        try:
            self._write_pdf(report_data, output_path)
        except Exception as error:
            logger.error("Failed to generate PDF report: %s", error, exc_info=True)
            raise HTTPException(status_code=500, detail=f"Failed to generate PDF report: {error}") from error

        llm_insights = report_data.get("llm_insights", {})
        metadata = {
            "format": "pdf",
            "file_path": str(output_path),
            "title": report_data["title"],
            "video_name": report_data["video_name"],
            "created_at": report_data["created_at"],
            "sections": 9,
            "section_names": [
                "Title + Metadata",
                "TL;DR",
                "Executive Summary",
                "Key Topics",
                "Scene Analysis",
                "Object Detection Results",
                "Graph Analysis",
                "Key Insights",
                "Conclusion",
            ],
            "content": {
                "objects_detected": len(report_data.get("detected_objects", [])),
                "key_topics": len(llm_insights.get("key_topics", [])),
                "key_insights": len(llm_insights.get("key_insights", [])),
                "graphs_detected": report_data.get("vision_metadata", {}).get("graphs_detected", False),
            },
            "llm_mode": "LLM-generated" if llm_insights else "heuristic",
        }

        response_text = self._format_report_response(report_data, output_path, "pdf")
        return AgentResult(
            agent=self.name,
            action="generate_pdf",
            response=response_text,
            confidence=1.0,
            metadata=metadata,
        )

    def generate_ppt(self, context: AgentContext) -> AgentResult:
        report_data = self._collect_analysis_inputs(context)
        report_dir = self._report_directory(context)
        output_path = report_dir / f"{self._safe_stem(context, 'presentation')}.pptx"

        try:
            self._write_ppt(report_data, output_path)
        except Exception as error:
            logger.error("Failed to generate PPT report: %s", error, exc_info=True)
            raise HTTPException(status_code=500, detail=f"Failed to generate PPT report: {error}") from error

        llm_insights = report_data.get("llm_insights", {})
        metadata = {
            "format": "ppt",
            "file_path": str(output_path),
            "title": report_data["title"],
            "video_name": report_data["video_name"],
            "created_at": report_data["created_at"],
            "slides": 6,
            "slide_names": [
                "Title",
                "Executive Summary",
                "Key Topics",
                "Objects Detected",
                "Key Insights",
                "Conclusion",
            ],
            "content": {
                "key_topics": len(llm_insights.get("key_topics", [])),
                "key_insights": len(llm_insights.get("key_insights", [])),
                "objects_detected": len(report_data.get("detected_objects", [])),
            },
            "llm_mode": "LLM-generated" if llm_insights else "heuristic",
        }

        response_text = self._format_report_response(report_data, output_path, "ppt")
        return AgentResult(
            agent=self.name,
            action="generate_ppt",
            response=response_text,
            confidence=1.0,
            metadata=metadata,
        )
