from io import BytesIO
from collections import Counter
from typing import Any

from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

from schemas import ExportRequest


def _metadata_value(payload: ExportRequest, key: str, default: Any = None) -> Any:
    if key in payload.metadata:
        return payload.metadata.get(key, default)
    return getattr(payload, key, default)


def _as_string_list(value: Any) -> list[str]:
    if value is None:
        return []
    if isinstance(value, list):
        return [str(item).strip() for item in value if str(item).strip()]
    if isinstance(value, tuple):
        return [str(item).strip() for item in value if str(item).strip()]
    if isinstance(value, str):
        return [value.strip()] if value.strip() else []
    return [str(value).strip()] if str(value).strip() else []


def _summarize_text(text: str, max_sentences: int = 4) -> str:
    cleaned = (text or "").strip()
    if not cleaned:
        return "No summary was available from the source analysis."

    sentences = [sentence.strip() for sentence in cleaned.replace("\n", " ").split(".") if sentence.strip()]
    if not sentences:
        return cleaned[:600]

    summary = ". ".join(sentences[:max_sentences]).strip()
    if summary and not summary.endswith("."):
        summary += "."
    return summary


def _top_items(items: list[str], limit: int = 6) -> list[str]:
    cleaned = [item.strip() for item in items if item and item.strip()]
    if not cleaned:
        return []
    counts = Counter(cleaned)
    return [item for item, _ in counts.most_common(limit)]


def _vision_insight_lines(vision_summary: Any) -> list[str]:
    if not isinstance(vision_summary, dict):
        return []

    insights: list[str] = []
    for key, value in vision_summary.items():
        if isinstance(value, (str, int, float, bool)):
            insights.append(f"{key}: {value}")
        elif isinstance(value, list) and value:
            insights.append(f"{key}: {len(value)} observations")
    return insights


def write_wrapped_pdf_text(pdf: canvas.Canvas, text: str, x: int, y: int, max_width: int, line_height: int) -> int:
    for paragraph in text.split("\n"):
        words = paragraph.split()

        if not words:
            y -= line_height
            continue

        current_line = words[0]
        for word in words[1:]:
            candidate_line = f"{current_line} {word}"
            if pdf.stringWidth(candidate_line, "Helvetica", 11) <= max_width:
                current_line = candidate_line
            else:
                pdf.drawString(x, y, current_line)
                y -= line_height
                current_line = word

        pdf.drawString(x, y, current_line)
        y -= line_height

    return y


def generate_pdf_report(payload: ExportRequest) -> BytesIO:
    buffer = BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=A4)
    page_width, page_height = A4
    y_position = page_height - 60

    transcript_summary = _metadata_value(payload, "transcript_summary", payload.transcription)
    detected_objects = _as_string_list(_metadata_value(payload, "detected_objects", payload.detected_objects))
    frame_text = _as_string_list(_metadata_value(payload, "frame_text", payload.frame_text))
    vision_summary = _metadata_value(payload, "vision_summary", {}) or {}

    executive_summary = _summarize_text(transcript_summary or payload.transcription)
    top_objects = _top_items(detected_objects, limit=6)
    top_text_snippets = _top_items(frame_text, limit=6)
    vision_insights = _vision_insight_lines(vision_summary)

    pdf.setTitle(payload.title)
    pdf.setFont("Helvetica-Bold", 18)
    pdf.drawString(50, y_position, payload.title)
    y_position -= 26

    sections = [
        (
            "Executive Summary",
            executive_summary,
        ),
        (
            "Key Context Insights",
            "\n".join(
                [
                    "- Visual focus: " + (", ".join(top_objects) if top_objects else "No dominant objects were detected."),
                    "- On-screen themes: " + (" | ".join(top_text_snippets) if top_text_snippets else "No clear text themes were extracted."),
                    "- Supporting signal count: "
                    + f"{len(detected_objects)} object cues and {len(frame_text)} text cues were identified.",
                ]
            ),
        ),
        (
            "Interpretation",
            "\n".join(
                [
                    "The video appears to emphasize a small set of recurring ideas, supported by repeated visual and textual cues.",
                    "These cues can be used to shape a concise narrative of the video's message and intended audience takeaway.",
                    *( ["Supplementary analysis: " + "; ".join(vision_insights)] if vision_insights else [] ),
                ]
            ),
        ),
        (
            "Recommended Next Steps",
            "\n".join(
                [
                    "- Use the executive summary as the opening context in presentations or documentation.",
                    "- Validate the most important insights against key timestamps in the original video.",
                    "- Convert the top context insights into action items for the target audience.",
                ]
            ),
        ),
    ]

    for heading, content in sections:
        if y_position < 120:
            pdf.showPage()
            y_position = page_height - 60

        pdf.setFont("Helvetica-Bold", 13)
        pdf.drawString(50, y_position, heading)
        y_position -= 20
        pdf.setFont("Helvetica", 11)
        y_position = write_wrapped_pdf_text(pdf, content, 50, y_position, int(page_width - 100), 16)
        y_position -= 10

    pdf.save()
    buffer.seek(0)
    return buffer


def add_ppt_body_slide(presentation: Presentation, title: str, body: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    chunks = [chunk.strip() for chunk in body.split("\n") if chunk.strip()]
    if not chunks:
        chunks = ["No content available."]

    for index, chunk in enumerate(chunks):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = chunk


def generate_ppt_report(payload: ExportRequest) -> BytesIO:
    presentation = Presentation()

    transcript_summary = _metadata_value(payload, "transcript_summary", payload.transcription)
    detected_objects = _as_string_list(_metadata_value(payload, "detected_objects", payload.detected_objects))
    frame_text = _as_string_list(_metadata_value(payload, "frame_text", payload.frame_text))
    vision_summary = _metadata_value(payload, "vision_summary", {}) or {}

    cover_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    cover_slide.shapes.title.text = payload.title
    cover_slide.placeholders[1].text = (
        f"Video: {payload.filename or 'Unknown'}\n"
        f"Agent: {payload.agent}\n"
        f"Language: {payload.language}\n"
        f"Duration: {payload.duration}"
    )

    add_ppt_body_slide(presentation, "Query", payload.query)
    add_ppt_body_slide(presentation, "Key Points", transcript_summary or "No transcript summary available.")
    add_ppt_body_slide(
        presentation,
        "Findings",
        "\n".join(
            [
                "Objects: " + (", ".join(detected_objects) if detected_objects else "none"),
                "Frame text: " + (" | ".join(frame_text[:5]) if frame_text else "none"),
                "Vision insights: " + (str(vision_summary) if vision_summary else "none"),
            ]
        ),
    )
    add_ppt_body_slide(
        presentation,
        "Conclusion",
        payload.response,
    )

    buffer = BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer


def generate_export_file(payload: ExportRequest) -> tuple[BytesIO, str, str]:
    export_format = payload.format.lower().strip()
    safe_stem = (payload.filename or payload.title or "video-ai-report").replace(" ", "-").lower()

    if export_format == "pdf":
        return generate_pdf_report(payload), "application/pdf", f"{safe_stem}.pdf"

    if export_format == "ppt":
        return (
            generate_ppt_report(payload),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            f"{safe_stem}.pptx",
        )

    raise ValueError("Unsupported export format. Use 'pdf' or 'ppt'.")